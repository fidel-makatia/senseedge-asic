#!/usr/bin/env python3
"""Generate SenseEdge demo presentation for 3-minute video."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

DOCS = os.path.dirname(os.path.abspath(__file__))
PROJECT = os.path.dirname(DOCS)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors
DARK_BG = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT_BLUE = RGBColor(0x4C, 0x6E, 0xF5)
ACCENT_GREEN = RGBColor(0x2E, 0xCC, 0x71)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
DARK_TEXT = RGBColor(0x2D, 0x2D, 0x2D)
ORANGE = RGBColor(0xF3, 0x96, 0x21)

def add_dark_bg(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG

def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Arial"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_bullet_slide(slide, left, top, width, height, items, font_size=16, color=WHITE):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Arial"
        p.space_after = Pt(8)
    return txBox

def add_image_safe(slide, path, left, top, width=None, height=None):
    if os.path.exists(path):
        kwargs = {}
        if width: kwargs['width'] = Inches(width)
        if height: kwargs['height'] = Inches(height)
        slide.shapes.add_picture(path, Inches(left), Inches(top), **kwargs)
        return True
    return False

def add_accent_bar(slide, left, top, width, height, color=ACCENT_BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top),
                                    Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

# ============================================================
# SLIDE 1: Title
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
add_dark_bg(slide)
add_accent_bar(slide, 0, 0, 13.333, 0.08, ACCENT_BLUE)
add_accent_bar(slide, 0, 7.42, 13.333, 0.08, ACCENT_BLUE)

add_text_box(slide, 1, 1.5, 11, 1.2, "SenseEdge", 54, ACCENT_BLUE, True, PP_ALIGN.CENTER)
add_text_box(slide, 1, 2.7, 11, 0.8, "Predictive Maintenance ASIC with Hardware FFT & Neural Network",
             24, WHITE, False, PP_ALIGN.CENTER)
add_text_box(slide, 1, 3.8, 11, 0.5, "Open-Source Custom Silicon on SkyWater SKY130",
             18, LIGHT_GRAY, False, PP_ALIGN.CENTER)

add_accent_bar(slide, 5.5, 4.5, 2.3, 0.04, ACCENT_BLUE)

add_text_box(slide, 1, 5.0, 11, 0.5, "Fidel Makatia", 22, WHITE, True, PP_ALIGN.CENTER)
add_text_box(slide, 1, 5.5, 11, 0.5, "PhD Student, Analog & Mixed-Signal IC Design", 16, LIGHT_GRAY, False, PP_ALIGN.CENTER)
add_text_box(slide, 1, 6.2, 11, 0.5, "ChipFoundry Reference Application Design Contest 2026", 14, ACCENT_BLUE, False, PP_ALIGN.CENTER)

# ============================================================
# SLIDE 2: Problem
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)
add_accent_bar(slide, 0, 0, 13.333, 0.08, ORANGE)

add_text_box(slide, 0.8, 0.4, 11, 0.8, "The Problem", 36, ORANGE, True)
add_accent_bar(slide, 0.8, 1.1, 2, 0.04, ORANGE)

add_bullet_slide(slide, 0.8, 1.5, 5.5, 5, [
    "Unplanned equipment downtime costs $50B/year",
    "",
    "Commercial vibration sensors: $500 - $5,000 per node",
    "",
    "Cloud-dependent ML adds latency & recurring costs",
    "",
    "Fixed-threshold alerting = excessive false alarms",
    "",
    "90% of factories can't afford predictive maintenance",
], 18, WHITE)

# Stats box
add_accent_bar(slide, 7, 1.5, 5.5, 4.5, RGBColor(0x25, 0x25, 0x40))
add_text_box(slide, 7.3, 1.8, 5, 0.5, "Market Opportunity", 20, ACCENT_BLUE, True, PP_ALIGN.CENTER)
add_text_box(slide, 7.3, 2.5, 5, 0.5, "$50B", 48, ORANGE, True, PP_ALIGN.CENTER)
add_text_box(slide, 7.3, 3.3, 5, 0.5, "Annual downtime cost", 16, LIGHT_GRAY, False, PP_ALIGN.CENTER)
add_text_box(slide, 7.3, 4.0, 5, 0.5, "$30B", 48, ACCENT_GREEN, True, PP_ALIGN.CENTER)
add_text_box(slide, 7.3, 4.8, 5, 0.5, "Predictive maintenance market by 2028", 16, LIGHT_GRAY, False, PP_ALIGN.CENTER)

# ============================================================
# SLIDE 3: Solution
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)
add_accent_bar(slide, 0, 0, 13.333, 0.08, ACCENT_GREEN)

add_text_box(slide, 0.8, 0.4, 11, 0.8, "The Solution: Full Pipeline in Silicon", 36, ACCENT_GREEN, True)
add_accent_bar(slide, 0.8, 1.1, 3, 0.04, ACCENT_GREEN)

# Pipeline flow
boxes = [
    ("SPI ADC\nInterface", "12-bit\n100 kSPS"),
    ("64-Point\nFFT Engine", "Radix-2 DIT\n24-bit internal"),
    ("Feature\nExtraction", "8 spectral\nfeatures"),
    ("Neural Net\nInference", "8→16→4 FC\nINT8, ReLU"),
    ("Alarm\nLogic", "GPIO + IRQ\nto RISC-V"),
]

for i, (title, desc) in enumerate(boxes):
    x = 0.5 + i * 2.5
    add_accent_bar(slide, x, 2.0, 2.2, 1.8, RGBColor(0x25, 0x25, 0x40))
    add_accent_bar(slide, x, 2.0, 2.2, 0.06, ACCENT_BLUE)
    add_text_box(slide, x + 0.1, 2.15, 2, 0.8, title, 14, ACCENT_BLUE, True, PP_ALIGN.CENTER)
    add_text_box(slide, x + 0.1, 3.0, 2, 0.6, desc, 11, LIGHT_GRAY, False, PP_ALIGN.CENTER)
    if i < 4:
        add_text_box(slide, x + 2.15, 2.6, 0.4, 0.5, "→", 24, ACCENT_BLUE, True, PP_ALIGN.CENTER)

# Key stats
add_bullet_slide(slide, 0.8, 4.3, 5.5, 3, [
    "< 10 μs inference latency (192 MACs at 20 MHz)",
    "Zero CPU intervention — fully autonomous pipeline",
    "212 runtime-loadable INT8 weights (field-updateable)",
    "< 5 mW total power at 1.8V",
], 15, WHITE)

add_accent_bar(slide, 7, 4.3, 5.5, 2.5, RGBColor(0x25, 0x25, 0x40))
add_text_box(slide, 7.3, 4.5, 5, 0.5, "Cost Comparison", 18, ACCENT_BLUE, True, PP_ALIGN.CENTER)
add_text_box(slide, 7.3, 5.0, 2.2, 0.8, "Commercial\n$500-$5,000", 16, RGBColor(0xFF, 0x55, 0x55), True, PP_ALIGN.CENTER)
add_text_box(slide, 9.5, 5.0, 0.5, 0.5, "vs", 16, LIGHT_GRAY, False, PP_ALIGN.CENTER)
add_text_box(slide, 10, 5.0, 2.2, 0.8, "SenseEdge\n~$13", 16, ACCENT_GREEN, True, PP_ALIGN.CENTER)
add_text_box(slide, 7.3, 6.0, 5, 0.5, "20-300x cost reduction", 14, ORANGE, True, PP_ALIGN.CENTER)

# ============================================================
# SLIDE 4: Silicon Design & Verification
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)
add_accent_bar(slide, 0, 0, 13.333, 0.08, ACCENT_BLUE)

add_text_box(slide, 0.8, 0.4, 11, 0.8, "Silicon Design & Verification", 36, ACCENT_BLUE, True)
add_accent_bar(slide, 0.8, 1.1, 3.5, 0.04, ACCENT_BLUE)

# Left: specs
add_accent_bar(slide, 0.8, 1.5, 5.5, 2.8, RGBColor(0x25, 0x25, 0x40))
add_text_box(slide, 1, 1.6, 5, 0.4, "Hardening Results (LibreLane 2.4.6)", 16, ACCENT_BLUE, True)
add_bullet_slide(slide, 1, 2.1, 5, 2, [
    "Process: SkyWater SKY130 (130nm)",
    "Gate count: 44,409",
    "Die area: 2920 × 2500 μm",
    "Clock: 20 MHz (50ns period)",
    "DRC: CLEAN | LVS: CLEAN",
    "Setup slack (typical): +6.66 ns PASS",
], 13, WHITE)

# Right: verification
add_accent_bar(slide, 7, 1.5, 5.5, 2.8, RGBColor(0x25, 0x25, 0x40))
add_text_box(slide, 7.2, 1.6, 5, 0.4, "Verification (46 assertions, 0 failures)", 16, ACCENT_GREEN, True)
tests = [
    "SPI ADC Interface — PASS",
    "64-Point FFT Engine — PASS",
    "Feature Extraction — PASS",
    "Neural Network Engine — PASS",
    "Alarm Logic — PASS",
    "Wishbone Interface — PASS",
    "Full Integration — PASS",
]
add_bullet_slide(slide, 7.2, 2.1, 5, 2, tests, 13, WHITE)

# GL sim results image
add_image_safe(slide, os.path.join(DOCS, "gl_simulation_results.png"), 0.8, 4.5, 12, 2.8)

# ============================================================
# SLIDE 5: PCB Design
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)
add_accent_bar(slide, 0, 0, 13.333, 0.08, ACCENT_BLUE)

add_text_box(slide, 0.8, 0.4, 11, 0.8, "PCB Design", 36, ACCENT_BLUE, True)
add_accent_bar(slide, 0.8, 1.1, 1.5, 0.04, ACCENT_BLUE)

# PCB 3D render
add_image_safe(slide, os.path.join(DOCS, "pcb_kicad_3D.png"), 0.5, 1.5, 6, 3.5)

# Component list
add_accent_bar(slide, 7, 1.5, 5.5, 3.5, RGBColor(0x25, 0x25, 0x40))
add_text_box(slide, 7.2, 1.6, 5, 0.4, "Key Components", 18, ACCENT_BLUE, True)
add_bullet_slide(slide, 7.2, 2.1, 5, 2.5, [
    "U3: Caravel ASIC (QFN-64)",
    "U2: ADXL326 3-axis accelerometer",
    "J1: MCP3201 12-bit SPI ADC",
    "U1: ESP32-C3 WiFi/BLE",
    "U4: W25Q32 SPI Flash",
    "Y1: 20 MHz CMOS oscillator",
    "J2: USB-C power input",
], 14, WHITE)

# Schematic
add_image_safe(slide, os.path.join(DOCS, "kicad_schema.png"), 0.5, 5.2, 6, 2)

# Board specs
add_accent_bar(slide, 7, 5.2, 5.5, 2, RGBColor(0x25, 0x25, 0x40))
add_text_box(slide, 7.2, 5.3, 5, 0.4, "Board Specifications", 16, ACCENT_BLUE, True)
add_bullet_slide(slide, 7.2, 5.8, 5, 1.5, [
    "65mm × 45mm, 2-layer FR4",
    "Total BOM: ~$13.15 at 100 units",
    "Gerbers ready for fabrication",
    "USB-C powered (5V input)",
], 13, WHITE)

# ============================================================
# SLIDE 6: Enclosure
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)
add_accent_bar(slide, 0, 0, 13.333, 0.08, ACCENT_BLUE)

add_text_box(slide, 0.8, 0.4, 11, 0.8, "Mechanical Enclosure", 36, ACCENT_BLUE, True)
add_accent_bar(slide, 0.8, 1.1, 2.5, 0.04, ACCENT_BLUE)

# Closed enclosure
add_image_safe(slide, os.path.join(DOCS, "enclosure.png"), 0.3, 1.5, 6, 4)

# Open with PCB
add_image_safe(slide, os.path.join(DOCS, "enclosure_withpcb.png"), 6.8, 1.5, 6, 4)

# Specs
add_accent_bar(slide, 0.8, 5.8, 11.7, 1.4, RGBColor(0x25, 0x25, 0x40))
add_bullet_slide(slide, 1, 5.9, 5, 1.2, [
    "69.5 × 49.5 × 20mm — snap-fit design",
    "IP54-rated dust & splash protection",
    "3D printed ASA/PETG — industrial temp range",
], 14, WHITE)
add_bullet_slide(slide, 7, 5.9, 5, 1.2, [
    "M3 mounting ears for equipment attachment",
    "USB-C port cutout aligned to PCB",
    "LED light pipe, ventilation slots",
], 14, WHITE)

# ============================================================
# SLIDE 7: Applications & Impact
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)
add_accent_bar(slide, 0, 0, 13.333, 0.08, ACCENT_GREEN)

add_text_box(slide, 0.8, 0.4, 11, 0.8, "Applications & Impact", 36, ACCENT_GREEN, True)
add_accent_bar(slide, 0.8, 1.1, 2.5, 0.04, ACCENT_GREEN)

apps = [
    ("Manufacturing", "Motor, pump, compressor monitoring"),
    ("HVAC", "Fan & compressor health in buildings"),
    ("Energy", "Wind turbine gearbox monitoring"),
    ("Structural Health", "Bridge & infrastructure vibration"),
    ("Transportation", "Rail wheel & axle bearing diagnostics"),
    ("Oil & Gas", "Downhole pump monitoring"),
]

for i, (sector, desc) in enumerate(apps):
    row = i // 3
    col = i % 3
    x = 0.8 + col * 4.2
    y = 1.5 + row * 1.8
    add_accent_bar(slide, x, y, 3.8, 1.5, RGBColor(0x25, 0x25, 0x40))
    add_accent_bar(slide, x, y, 3.8, 0.05, ACCENT_BLUE)
    add_text_box(slide, x + 0.2, y + 0.15, 3.4, 0.5, sector, 18, ACCENT_BLUE, True)
    add_text_box(slide, x + 0.2, y + 0.7, 3.4, 0.6, desc, 13, LIGHT_GRAY)

# Key enabler
add_accent_bar(slide, 0.8, 5.5, 11.7, 1.5, RGBColor(0x25, 0x25, 0x40))
add_text_box(slide, 1, 5.6, 11, 0.4, "One Silicon, Every Application", 20, ORANGE, True, PP_ALIGN.CENTER)
add_text_box(slide, 1, 6.1, 11, 0.8,
             "Field-updateable neural network — 212 bytes of new weights transforms the sensor "
             "from motor monitoring to bridge health to HVAC diagnostics. No hardware changes, no re-tapeout.",
             14, LIGHT_GRAY, False, PP_ALIGN.CENTER)

# ============================================================
# SLIDE 8: Summary & Thank You
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)
add_accent_bar(slide, 0, 0, 13.333, 0.08, ACCENT_BLUE)
add_accent_bar(slide, 0, 7.42, 13.333, 0.08, ACCENT_BLUE)

add_text_box(slide, 1, 0.8, 11, 0.8, "SenseEdge — Complete Reference Design", 36, ACCENT_BLUE, True, PP_ALIGN.CENTER)
add_accent_bar(slide, 5.5, 1.5, 2.3, 0.04, ACCENT_BLUE)

# Deliverables grid
deliverables = [
    ("Custom ASIC", "44K gates, DRC/LVS clean\nSKY130, 20 MHz"),
    ("Verification", "7 testbenches, 46 assertions\nGL simulation confirmed"),
    ("PCB Design", "65×45mm, 2-layer FR4\nGerbers ready, ~$13 BOM"),
    ("Enclosure", "3D printable, IP54\nSnap-fit, USB-C cutout"),
    ("Firmware", "RISC-V C firmware\nUART reporting, weight loading"),
    ("ML Pipeline", "TensorFlow QAT\n212 INT8 parameters"),
]

for i, (title, desc) in enumerate(deliverables):
    row = i // 3
    col = i % 3
    x = 1 + col * 3.9
    y = 2.0 + row * 1.8
    add_accent_bar(slide, x, y, 3.5, 1.5, RGBColor(0x25, 0x25, 0x40))
    add_accent_bar(slide, x, y, 3.5, 0.05, ACCENT_GREEN)
    add_text_box(slide, x + 0.2, y + 0.15, 3.1, 0.5, title, 16, ACCENT_GREEN, True, PP_ALIGN.CENTER)
    add_text_box(slide, x + 0.2, y + 0.6, 3.1, 0.8, desc, 12, LIGHT_GRAY, False, PP_ALIGN.CENTER)

add_text_box(slide, 1, 5.8, 11, 0.5, "Thank You", 28, WHITE, True, PP_ALIGN.CENTER)
add_text_box(slide, 1, 6.3, 11, 0.4, "Fidel Makatia  |  ChipFoundry Reference Application Design Contest 2026", 14, LIGHT_GRAY, False, PP_ALIGN.CENTER)
add_text_box(slide, 1, 6.7, 11, 0.4, "github.com/fidel-makatia/senseedge-asic", 14, ACCENT_BLUE, False, PP_ALIGN.CENTER)

# ============================================================
# Save
# ============================================================
output_path = os.path.join(DOCS, "SenseEdge_Presentation.pptx")
prs.save(output_path)
print(f"Presentation saved: {output_path}")
print(f"Slides: {len(prs.slides)}")
